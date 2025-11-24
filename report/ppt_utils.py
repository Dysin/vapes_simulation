'''
@Desc:   输出PPT报告
@Author: Dysin
@Date:   2025/10/22
'''

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


class PptUtils:
    """
    PowerPoint 工具类
    用于快速创建幻灯片、添加标题、文本、图片等内容。
    """

    def __init__(self, ppt_path=None):
        """
        初始化 PPT 文件
        :param ppt_path: 现有PPT路径；若为None则新建
        """
        if ppt_path and os.path.exists(ppt_path):
            self.prs = Presentation(ppt_path)
        else:
            self.prs = Presentation()

        # 默认幻灯片尺寸：16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title: str, subtitle: str = ""):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
        return slide

    def add_text_slide(self, title: str, content: str):
        """添加带文本内容的页"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        return slide

    def add_blank_slide(self):
        """添加空白页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        return slide

    def add_picture(self, slide, image_path: str, left=1, top=1, width=None, height=None):
        """在指定幻灯片添加图片"""
        left = Inches(left)
        top = Inches(top)
        width = Inches(width) if width else None
        height = Inches(height) if height else None
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    def add_textbox(self, slide, text: str, left=1, top=1, width=4, height=1,
                    font_size=20, bold=False, color=(0, 0, 0), align="left"):
        """在幻灯片添加自定义文本框"""
        left, top, width, height = map(Inches, [left, top, width, height])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)

        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        return textbox

    def save(self, file_path: str):
        """保存PPT"""
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        self.prs.save(file_path)
        print(f"✅ PPT 已保存至：{file_path}")

if __name__ == '__main__':
    # 初始化
    ppt = PptUtils()

    # 添加标题页
    ppt.add_title_slide("电子烟气道仿真报告", "单相流模拟结果分析")

    # 添加文字页
    ppt.add_text_slide("仿真参数", "入口流速：15.6 m/s\n出口压力：101325 Pa")

    # 添加空白页并插图
    slide = ppt.add_blank_slide()
    ppt.add_picture(slide, "geometry_inlet.png", left=1, top=1, width=5)

    # 添加文本框
    ppt.add_textbox(slide, "速度分布云图", left=1, top=6, width=4, height=1,
                    font_size=18, bold=True, align="center")
    # 保存
    ppt.save("output/simulation_report.pptx")

